from pptx import Presentation

# Create a presentation object
ppt_updated = Presentation()

# Add a title slide
slide_title = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "Updated Marketing Campaign Data Analytics"
subtitle.text = "Refined Analysis and Results"

# Slide 1: Overview of the Project
slide1 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title1 = slide1.shapes.title
title1.text = "Project Overview"
content1 = slide1.shapes.placeholders[1].text = """
- Analysis of traditional and online marketing campaigns.
- Data-driven insights into campaign effectiveness.
- Predictive modeling to forecast campaign success or failure.
"""

# Slide 2: Data Sources
slide2 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "Data Sources"
content2 = slide2.shapes.placeholders[1].text = """
- Traditional Campaign Dataset:
  - Sales data, Market Size, Store Age, Promotions.
- Online Campaign Dataset (Facebook):
  - Metrics including Impressions, Clicks, Cost-per-Click (CPC), etc.
"""

# Slide 3: Data Cleaning and Preprocessing
slide3 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "Data Cleaning and Preprocessing"
content3 = slide3.shapes.placeholders[1].text = """
- Handled missing values in online campaign data.
- Applied one-hot encoding to convert categorical variables.
- Dropped original categorical columns post-encoding.
"""

# Slide 4: Exploratory Data Analysis (EDA)
slide4 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Exploratory Data Analysis (EDA)"
content4 = slide4.shapes.placeholders[1].text = """
- Visualized sales by promotion type using box plots.
- Examined impressions vs clicks in online campaigns using scatter plots.
"""

# Slide 5: Machine Learning Models
slide5 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Machine Learning Models"
content5 = slide5.shapes.placeholders[1].text = """
- Models used:
  - Random Forest for both Traditional and Online Campaigns.
- Predictions made for Sales (Traditional) and CTR (Online).
"""

# Slide 6: Results and Evaluation
slide6 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Results and Evaluation"
content6 = slide6.shapes.placeholders[1].text = """
- RMSE for Traditional Campaign: Calculated.
- RMSE for Online Campaign: Calculated.
- Visualization of Actual vs Predicted for both campaigns.
"""

# Slide 7: Future Work
slide7 = ppt_updated.slides.add_slide(ppt_updated.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Future Work"
content7 = slide7.shapes.placeholders[1].text = """
- Further model optimization using advanced algorithms.
- Incorporate additional variables (weather, social events).
- Compare effectiveness across multiple campaigns.
"""

# Save the updated presentation
ppt_updated_file = "Updated_Marketing_Campaign_Presentation.pptx"
ppt_updated.save(ppt_updated_file)

